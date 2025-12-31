"""PPTX parser for ParseForge."""

import logging
from pathlib import Path

from pptx import Presentation

from src.schema.document import BlockType, Document, Page, TextBlock, TitleBlock
from src.utils.exceptions import ParserError

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str) -> Document:
    """
    Parse PPTX file.

    Args:
        file_path: Path to PPTX file

    Returns:
        Document object
    """
    try:
        prs = Presentation(file_path)
        all_blocks = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Determine if title or text
                    if shape.placeholder_format.idx == 0:  # Title placeholder
                        block = TitleBlock(
                            block_type=BlockType.TITLE,
                            text=shape.text,
                            page_index=slide_idx,
                            level=1,
                        )
                    else:
                        block = TextBlock(
                            block_type=BlockType.TEXT,
                            text=shape.text,
                            page_index=slide_idx,
                        )
                    all_blocks.append(block)

        # Group blocks by page
        pages = []
        for slide_idx in range(len(prs.slides)):
            slide_blocks = [b for b in all_blocks if b.page_index == slide_idx]
            page = Page(
                page_index=slide_idx,
                width=960,  # Standard slide width
                height=540,  # Standard slide height
                blocks=slide_blocks,
            )
            pages.append(page)

        document = Document(
            file_path=file_path,
            file_name=Path(file_path).name,
            pages=pages,
            total_pages=len(pages),
        )

        return document

    except Exception as e:
        raise ParserError(f"Failed to parse PPTX: {e}") from e

